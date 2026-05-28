#!/usr/bin/env python3
"""
GAIA CodeAgent Python Step Runner — ADR-138 iter 54

This script is invoked by gaia-codeagent.ts for each agent step.
It pre-defines tool functions as Python callables, then exec()s the agent's
code in that context.

Protocol (via environment variables):
  GAIA_CODE_FILE     : path to a temp file containing the agent's Python code
  GAIA_RESULT_FILE   : path where this runner writes final_answer JSON sentinel
  ANTHROPIC_API_KEY  : passed through to claude -p tool backends
  GAIA_ATTACHMENT_PATH: optional attachment file path for this question

Tool dispatch strategy:
  web_search(query)       -> claude -p with WebSearch (best web coverage)
  visit_webpage(url)      -> requests + bs4 HTML extraction (no claude -p overhead)
  grounded_query(query)   -> ruflo TS grounded_query via node subprocess
  read_file(path)         -> direct Python (text/csv/json/xlsx/pptx) or requests+bs4 for PDF
  describe_image(path)    -> claude -p with --allowedTools Read (Claude vision)
  final_answer(answer)    -> writes GAIA_RESULT_FILE JSON and sys.exit(0)

Authorized imports (pre-loaded):
  math, re, json, datetime, collections, itertools, statistics,
  fractions, decimal, string, unicodedata, operator, sys, os, pathlib,
  urllib.parse, pandas (optional), numpy (optional), sympy (optional),
  requests (optional), bs4 (optional), openpyxl (optional)

Refs: ADR-138, HAL-DEEP-STUDY.md, #2156, iter 54
"""

import sys
import os
import json
import subprocess
import tempfile
import math
import re
import datetime
from datetime import date, timedelta
from collections import Counter, defaultdict, OrderedDict
import itertools
import statistics
import fractions
import decimal
import operator
import string
import unicodedata
from urllib.parse import urlparse, quote, unquote
from pathlib import Path

# Optional heavy imports — graceful fallback
try:
    import pandas as pd
    import numpy as np
except ImportError:
    pd = None
    np = None

try:
    import sympy
except ImportError:
    sympy = None

try:
    import requests
    from bs4 import BeautifulSoup
    HAS_REQUESTS = True
except ImportError:
    requests = None
    BeautifulSoup = None
    HAS_REQUESTS = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
GAIA_RESULT_FILE = os.environ.get("GAIA_RESULT_FILE", "")
GAIA_CODE_FILE = os.environ.get("GAIA_CODE_FILE", "")
GAIA_ATTACHMENT_PATH = os.environ.get("GAIA_ATTACHMENT_PATH", "")

CLAUDE_CMD = "claude"
TOOL_TIMEOUT = 30  # seconds per tool call


# ---------------------------------------------------------------------------
# Tool: web_search
# Dispatches to claude -p with WebSearch for best web coverage
# ---------------------------------------------------------------------------

def web_search(query: str) -> str:
    """Search the web and return relevant snippets."""
    if not query or not query.strip():
        return "[web_search: empty query]"
    try:
        result = subprocess.run(
            [
                CLAUDE_CMD, "-p",
                "--allowedTools", "WebSearch",
                "--output-format", "text",
                f"Search the web for: {query}\n\nReturn the top results with titles, URLs, and key snippets. Be concise.",
            ],
            capture_output=True,
            text=True,
            timeout=TOOL_TIMEOUT,
            env={**os.environ, "ANTHROPIC_API_KEY": ANTHROPIC_API_KEY},
        )
        output = result.stdout.strip()
        if output:
            return output[:4000]
        if result.stderr:
            return f"[web_search: {result.stderr.strip()[:200]}]"
        return "[web_search: no results]"
    except subprocess.TimeoutExpired:
        return f"[web_search: timed out after {TOOL_TIMEOUT}s]"
    except FileNotFoundError:
        # claude CLI not available — fall back to DuckDuckGo scraping
        return _ddg_fallback(query)
    except Exception as e:
        return f"[web_search error: {e}]"


def _ddg_fallback(query: str) -> str:
    """Fallback web search via DuckDuckGo HTML scraping."""
    if not HAS_REQUESTS:
        return "[web_search: requests not available]"
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; ruflo-gaia/1.0)"}
        resp = requests.post(
            "https://html.duckduckgo.com/html/",
            data={"q": query},
            headers=headers,
            timeout=15,
        )
        soup = BeautifulSoup(resp.text, "html.parser")
        results = []
        for r in soup.select(".result")[:5]:
            title_el = r.select_one(".result__title")
            snip_el = r.select_one(".result__snippet")
            title = title_el.get_text(strip=True) if title_el else ""
            snip = snip_el.get_text(strip=True) if snip_el else ""
            if title:
                results.append(f"{title}: {snip}")
        return "\n".join(results) if results else "[web_search: no results]"
    except Exception as e:
        return f"[web_search: {e}]"


# ---------------------------------------------------------------------------
# Tool: visit_webpage
# Fetches full page content via requests + bs4
# ---------------------------------------------------------------------------

def visit_webpage(url: str) -> str:
    """Fetch the full text content of a webpage."""
    if not url or not url.strip():
        return "[visit_webpage: empty URL]"
    url = url.strip()
    if not url.startswith("http"):
        url = "https://" + url
    if HAS_REQUESTS:
        return _visit_with_requests(url)
    else:
        return _visit_with_claude(url)


def _visit_with_requests(url: str) -> str:
    """Fetch page via requests + bs4 text extraction."""
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; ruflo-gaia/1.0)"}
        resp = requests.get(url, headers=headers, timeout=20, allow_redirects=True)
        if resp.status_code != 200:
            return f"[visit_webpage: HTTP {resp.status_code} for {url}]"

        soup = BeautifulSoup(resp.text, "html.parser")
        # Remove noise
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "noscript"]):
            tag.decompose()

        # Extract main content
        text = soup.get_text(separator="\n", strip=True)
        # Collapse blank lines
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        content = "\n".join(lines)

        if len(content) > 8000:
            content = content[:8000] + "\n[... truncated ...]"
        return content if content else "[visit_webpage: no text content extracted]"
    except requests.exceptions.Timeout:
        return f"[visit_webpage: timed out fetching {url}]"
    except Exception as e:
        return f"[visit_webpage error: {e}]"


def _visit_with_claude(url: str) -> str:
    """Fallback: fetch webpage via claude -p with WebFetch."""
    try:
        result = subprocess.run(
            [
                CLAUDE_CMD, "-p",
                "--allowedTools", "WebFetch",
                "--output-format", "text",
                f"Fetch the content of this URL and return the full text: {url}",
            ],
            capture_output=True,
            text=True,
            timeout=TOOL_TIMEOUT,
            env={**os.environ, "ANTHROPIC_API_KEY": ANTHROPIC_API_KEY},
        )
        output = result.stdout.strip()
        return output[:8000] if output else f"[visit_webpage: no content from {url}]"
    except Exception as e:
        return f"[visit_webpage: {e}]"


# ---------------------------------------------------------------------------
# Tool: grounded_query
# Calls Gemini with Google Search grounding (ruflo's unique advantage)
# ---------------------------------------------------------------------------

def grounded_query(query: str) -> str:
    """Ask Gemini with Google Search grounding — returns synthesized answer with citations."""
    if not query or not query.strip():
        return "[grounded_query: empty query]"

    google_key = os.environ.get("GOOGLE_AI_API_KEY", "")
    if not google_key:
        # Fall back to web_search if no Google key
        return web_search(query)

    try:
        import urllib.request
        import urllib.error

        url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"
        payload = {
            "contents": [{"parts": [{"text": query}]}],
            "tools": [{"google_search": {}}],
            "generationConfig": {"maxOutputTokens": 1024},
        }
        data = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(
            f"{url}?key={google_key}",
            data=data,
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=20) as resp:
            body = json.loads(resp.read().decode("utf-8"))

        # Extract text
        candidates = body.get("candidates", [])
        if not candidates:
            return "[grounded_query: no candidates]"
        parts = candidates[0].get("content", {}).get("parts", [])
        text = " ".join(p.get("text", "") for p in parts).strip()

        # Extract sources
        sources = []
        grounding = candidates[0].get("groundingMetadata", {})
        for chunk in grounding.get("groundingChunks", [])[:4]:
            web = chunk.get("web", {})
            title = web.get("title", "")
            src_uri = web.get("uri", "")
            if title and src_uri:
                sources.append(f"- {title}: {src_uri}")

        result = text
        if sources:
            result += "\n\nSources:\n" + "\n".join(sources)
        return result[:4000] if result else "[grounded_query: empty response]"

    except Exception as e:
        # Fall back to web_search
        return web_search(query)


# ---------------------------------------------------------------------------
# Tool: read_file
# Supports: txt, csv, json, xlsx, pptx, pdf (via pdfminer.six if available)
# ---------------------------------------------------------------------------

def read_file(path: str) -> str:
    """Read a file and return its text content."""
    if not path or not path.strip():
        return "[read_file: empty path]"
    path = path.strip()
    if not os.path.exists(path):
        return f"[read_file: file not found: {path}]"

    ext = Path(path).suffix.lower()

    # Text-based formats — direct read
    if ext in (".txt", ".md", ".log", ".xml", ".html", ".htm"):
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(50_000)
            return content if content else "[read_file: empty file]"
        except Exception as e:
            return f"[read_file: {e}]"

    # JSON
    if ext == ".json":
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return json.dumps(data, indent=2)[:8000]
        except Exception as e:
            return f"[read_file json: {e}]"

    # CSV
    if ext == ".csv":
        if pd is not None:
            try:
                df = pd.read_csv(path)
                return f"Shape: {df.shape}\nColumns: {list(df.columns)}\n\nFirst rows:\n{df.head(10).to_string()}"
            except Exception as e:
                return f"[read_file csv: {e}]"
        else:
            try:
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    return f.read(8000)
            except Exception as e:
                return f"[read_file csv: {e}]"

    # Excel
    if ext in (".xlsx", ".xls"):
        if pd is not None:
            try:
                df = pd.read_excel(path, sheet_name=None)
                parts = []
                for sheet_name, sheet_df in df.items():
                    parts.append(f"Sheet: {sheet_name}")
                    parts.append(sheet_df.head(20).to_string())
                return "\n\n".join(parts)[:8000]
            except Exception as e:
                return f"[read_file excel: {e}]"
        elif HAS_OPENPYXL:
            try:
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                parts = []
                for ws in wb.worksheets:
                    rows = []
                    for row in ws.iter_rows(max_row=20, values_only=True):
                        rows.append("\t".join(str(c) if c is not None else "" for c in row))
                    parts.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
                return "\n\n".join(parts)[:8000]
            except Exception as e:
                return f"[read_file excel: {e}]"
        return "[read_file: pandas/openpyxl not available for Excel]"

    # PPTX
    if ext == ".pptx":
        if HAS_PPTX:
            try:
                prs = PptxPresentation(path)
                parts = []
                for i, slide in enumerate(prs.slides):
                    texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
                    parts.append(f"Slide {i+1}: " + " | ".join(texts))
                return "\n".join(parts)[:8000]
            except Exception as e:
                return f"[read_file pptx: {e}]"
        return "[read_file: python-pptx not available]"

    # PDF — try pdfminer.six, then fallback to pdftotext subprocess
    if ext == ".pdf":
        return _read_pdf(path)

    # Images — use describe_image
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"):
        return describe_image(path)

    # Unknown — try as text
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(8000)
    except Exception as e:
        return f"[read_file: unsupported format {ext}: {e}]"


def _read_pdf(path: str) -> str:
    """Extract text from a PDF file."""
    # Try pdfminer.six
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        text = pdfminer_extract(path)
        if text and text.strip():
            return text.strip()[:8000]
    except ImportError:
        pass
    except Exception:
        pass

    # Try pdftotext subprocess
    try:
        result = subprocess.run(
            ["pdftotext", path, "-"],
            capture_output=True, text=True, timeout=15
        )
        if result.stdout.strip():
            return result.stdout.strip()[:8000]
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception):
        pass

    # Last resort: claude -p with Read (vision for scanned PDFs)
    if ANTHROPIC_API_KEY:
        try:
            result = subprocess.run(
                [
                    CLAUDE_CMD, "-p",
                    "--allowedTools", "Read",
                    "--output-format", "text",
                    f"Read and extract all text from this file: {path}",
                ],
                capture_output=True, text=True, timeout=30,
                env={**os.environ, "ANTHROPIC_API_KEY": ANTHROPIC_API_KEY},
            )
            if result.stdout.strip():
                return result.stdout.strip()[:8000]
        except Exception:
            pass

    return f"[read_file: could not extract text from PDF at {path}]"


# ---------------------------------------------------------------------------
# Tool: describe_image
# Uses claude -p with vision capability
# ---------------------------------------------------------------------------

def describe_image(path: str) -> str:
    """Describe the contents of an image using Claude vision."""
    if not path or not os.path.exists(path):
        return f"[describe_image: file not found: {path}]"
    if not ANTHROPIC_API_KEY:
        return "[describe_image: ANTHROPIC_API_KEY not set]"
    try:
        result = subprocess.run(
            [
                CLAUDE_CMD, "-p",
                "--allowedTools", "Read",
                "--output-format", "text",
                f"Describe this image in detail, including any text, numbers, charts, or data visible: {path}",
            ],
            capture_output=True, text=True, timeout=30,
            env={**os.environ, "ANTHROPIC_API_KEY": ANTHROPIC_API_KEY},
        )
        output = result.stdout.strip()
        return output[:4000] if output else "[describe_image: no description generated]"
    except Exception as e:
        return f"[describe_image: {e}]"


# ---------------------------------------------------------------------------
# Tool: final_answer
# Writes sentinel JSON and exits
# ---------------------------------------------------------------------------

def final_answer(answer) -> None:
    """Submit the final answer. This ends the agent loop."""
    answer_str = str(answer).strip()

    # Write sentinel JSON
    if GAIA_RESULT_FILE:
        try:
            with open(GAIA_RESULT_FILE, "w", encoding="utf-8") as f:
                json.dump({"answer": answer_str}, f)
        except Exception as e:
            print(f"[final_answer: could not write result file: {e}]", file=sys.stderr)

    print(f"[GAIA_FINAL_ANSWER]: {answer_str}")
    sys.exit(0)


# ---------------------------------------------------------------------------
# Main execution
# ---------------------------------------------------------------------------

def main():
    # Read agent code
    if not GAIA_CODE_FILE or not os.path.exists(GAIA_CODE_FILE):
        print("[CodeAgent runner: GAIA_CODE_FILE not set or not found]", file=sys.stderr)
        sys.exit(1)

    with open(GAIA_CODE_FILE, "r", encoding="utf-8") as f:
        agent_code = f.read()

    # Build execution globals with all tool functions pre-defined
    exec_globals = {
        "__builtins__": __builtins__,
        # Math / standard lib
        "math": math,
        "re": re,
        "json": json,
        "datetime": datetime,
        "date": date,
        "timedelta": timedelta,
        "Counter": Counter,
        "defaultdict": defaultdict,
        "OrderedDict": OrderedDict,
        "itertools": itertools,
        "statistics": statistics,
        "fractions": fractions,
        "decimal": decimal,
        "operator": operator,
        "string": string,
        "unicodedata": unicodedata,
        "urlparse": urlparse,
        "quote": quote,
        "unquote": unquote,
        "Path": Path,
        "os": os,
        "sys": sys,
        "pd": pd,
        "np": np,
        "sympy": sympy,
        "requests": requests,
        # Tool functions
        "web_search": web_search,
        "visit_webpage": visit_webpage,
        "grounded_query": grounded_query,
        "read_file": read_file,
        "describe_image": describe_image,
        "final_answer": final_answer,
        # Attachment path hint
        "ATTACHMENT_PATH": GAIA_ATTACHMENT_PATH,
    }

    # Execute
    try:
        exec(agent_code, exec_globals)
    except SystemExit:
        # final_answer() calls sys.exit(0) — expected
        raise
    except Exception as e:
        import traceback
        print(f"[CodeAgent execution error]\n{traceback.format_exc()}", file=sys.stderr)
        sys.exit(2)


if __name__ == "__main__":
    main()
